#!/usr/bin/env python3
"""
DIS Group × Skyscanner — Enterprise Presentation Generator
Creates a professional 12-slide PPTX (16:9, white, consulting-grade)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ─── DESIGN SYSTEM ──────────────────────────────────────────────────────────
C_NAVY      = RGBColor(14,  44,  84)
C_BLUE      = RGBColor(37,  99, 235)
C_TEXT      = RGBColor(17,  24,  39)
C_MUTED     = RGBColor(107, 114, 128)
C_LIGHT_BG  = RGBColor(247, 248, 250)
C_BORDER    = RGBColor(226, 232, 240)
C_WHITE     = RGBColor(255, 255, 255)
C_GREEN     = RGBColor(5,   150, 105)
C_PURPLE    = RGBColor(88,  28, 135)
C_TEAL      = RGBColor(2,   90,  70)
C_LIGHT_BLUE_BG = RGBColor(235, 242, 255)

FONT = "Calibri"
W = Inches(13.333)
H = Inches(7.5)
LM = Inches(0.75)
RM = Inches(0.75)
CONTENT_W = W - LM - RM
HEADER_H  = Inches(0.78)
FOOTER_Y  = Inches(7.08)
FOOTER_H  = Inches(0.42)


# ─── HELPERS ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_WHITE
    return slide


def no_border(shape):
    spPr = shape._element.spPr
    for ln in spPr.findall(qn('a:ln')):
        spPr.remove(ln)
    ln_elem = etree.SubElement(spPr, qn('a:ln'))
    etree.SubElement(ln_elem, qn('a:noFill'))


def rect(slide, l, t, w, h, fill=None, line_rgb=None, line_pt=0.75):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(line_pt)
    else:
        no_border(shp)
    return shp


def txt(slide, text, l, t, w, h, size=11, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    for line in text.split('\n'):
        if p.runs:
            p = tf.add_paragraph()
            p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def header(slide, title, num, total=12):
    rect(slide, 0, 0, W, HEADER_H, fill=C_NAVY)
    txt(slide, f"{num:02d}", Inches(0.2), Inches(0.14),
        Inches(0.45), Inches(0.52), size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, title, Inches(0.62), Inches(0.14),
        W - Inches(4.0), HEADER_H - Inches(0.08),
        size=19, bold=True, color=C_WHITE)
    txt(slide, "DIS GROUP  ×  SKYSCANNER   |   CONFIDENTIAL",
        W - Inches(3.6), Inches(0.26), Inches(3.4), Inches(0.32),
        size=8, color=RGBColor(160, 185, 220), align=PP_ALIGN.RIGHT)


def footer(slide, num, total=12):
    rect(slide, 0, FOOTER_Y, W, FOOTER_H, fill=RGBColor(246, 248, 250))
    rect(slide, 0, FOOTER_Y, W, Inches(0.025), fill=C_BORDER)
    txt(slide, "DIS Group × Skyscanner  |  Confidential — For Partnership Discussion Only  |  July 2026",
        LM, FOOTER_Y + Inches(0.07), W - LM - Inches(1.4), Inches(0.3),
        size=8, color=C_MUTED)
    txt(slide, f"{num} / {total}",
        W - Inches(1.15), FOOTER_Y + Inches(0.07), Inches(1.0), Inches(0.3),
        size=8, color=C_MUTED, align=PP_ALIGN.RIGHT)


def content_top():
    return HEADER_H + Inches(0.12)


def content_height():
    return FOOTER_Y - content_top() - Inches(0.08)


def section_label(slide, text, y):
    txt(slide, text, LM, y, CONTENT_W, Inches(0.28),
        size=8.5, bold=True, color=C_BLUE)


# ─── SLIDES ─────────────────────────────────────────────────────────────────

def s01_title(prs):
    slide = blank_slide(prs)
    panel_w = Inches(5.1)

    # Left navy panel
    rect(slide, 0, 0, panel_w, H, fill=C_NAVY)

    # DIS GROUP wordmark
    txt(slide, "DIS GROUP", Inches(0.5), Inches(1.6), Inches(4.2), Inches(1.0),
        size=40, bold=True, color=C_WHITE)

    # Accent rule
    rect(slide, Inches(0.5), Inches(2.72), Inches(1.4), Pt(3), fill=C_BLUE)

    # Tagline
    txt(slide, "The Intelligence Layer\nfor Corporate Travel",
        Inches(0.5), Inches(2.88), Inches(4.2), Inches(0.95),
        size=15.5, color=RGBColor(175, 200, 235))

    # Product name
    txt(slide, "Introducing DIS Travel",
        Inches(0.5), Inches(3.95), Inches(4.2), Inches(0.45),
        size=11, color=RGBColor(140, 170, 215), italic=True)

    # Right side
    rl = Inches(5.5)
    rw = Inches(7.4)

    txt(slide, "PREPARED FOR", rl, Inches(1.52), rw, Inches(0.3),
        size=8.5, bold=True, color=C_MUTED)
    txt(slide, "Skyscanner", rl, Inches(1.82), rw, Inches(0.7),
        size=30, bold=True, color=C_NAVY)
    txt(slide, "Partnerships Team", rl, Inches(2.47), rw, Inches(0.42),
        size=16, color=C_MUTED)

    rect(slide, rl, Inches(3.0), rw, Inches(0.025), fill=C_BORDER)

    txt(slide, "PARTNERSHIP PACKAGE CONTENTS",
        rl, Inches(3.18), rw, Inches(0.28),
        size=8.5, bold=True, color=C_MUTED)

    items = [
        "01.  Cover Letter",
        "02.  Executive One-Pager",
        "03.  Company & Product Presentation  (this document)",
        "04.  Technical Architecture Overview",
    ]
    for i, item in enumerate(items):
        txt(slide, item, rl, Inches(3.52) + Inches(i * 0.42), rw, Inches(0.38),
            size=11, color=C_TEXT)

    rect(slide, rl, Inches(5.75), rw, Inches(0.025), fill=C_BORDER)
    txt(slide, "July 2026   |   Confidential — For Partnership Discussion Only",
        rl, Inches(5.88), rw, Inches(0.32), size=10, color=C_MUTED)

    return slide


def s02_about(prs):
    slide = blank_slide(prs)
    header(slide, "DIS Group at a Glance", 2)
    footer(slide, 2)
    ct = content_top()

    txt(slide, "One thesis applied across three markets: build the intelligence layer for industries that make decisions with poor data.",
        LM, ct, CONTENT_W, Inches(0.32), size=11, italic=True, color=C_MUTED)

    cw = Inches(3.87)
    ch = Inches(2.28)
    cy = ct + Inches(0.42)
    gap = Inches(0.18)

    cards = [
        ("DIS TRAVEL", "AI-powered Travel Intelligence Platform",
         "Intelligence layer for corporate buyers and travel-intensive organisations. Built for CSRD compliance from the foundation.",
         C_NAVY),
        ("DIS ENERGY", "Trading and Asset Decision Intelligence",
         "Intelligence layer for energy market decisions, asset portfolio management, and regulatory positioning.",
         C_TEAL),
        ("DIS ESG", "Environmental and Governance Data",
         "Structured ESG data infrastructure for corporate compliance, CSRD reporting, and capital markets disclosure.",
         C_PURPLE),
    ]
    for i, (name, sub, body, col) in enumerate(cards):
        x = LM + i * (cw + gap)
        rect(slide, x, cy, cw, ch, fill=col)
        txt(slide, name, x + Inches(0.18), cy + Inches(0.14), cw - Inches(0.36), Inches(0.38),
            size=11, bold=True, color=C_WHITE)
        txt(slide, sub, x + Inches(0.18), cy + Inches(0.50), cw - Inches(0.36), Inches(0.4),
            size=9.5, color=RGBColor(210, 225, 245), italic=True)
        txt(slide, body, x + Inches(0.18), cy + Inches(0.92), cw - Inches(0.36), ch - Inches(1.06),
            size=10, color=RGBColor(220, 232, 250))

    # Shared characteristics
    sy = cy + ch + Inches(0.22)
    section_label(slide, "WHAT EVERY DIS PRODUCT SHARES", sy)
    sy += Inches(0.3)

    shared = [
        "AI-native from inception — no legacy architecture to retrofit",
        "Data treated as a strategic asset: owned, structured, and governed",
        "GDPR-compliant design — EU-incorporated company, EU enterprise market",
        "Long-term horizon — built for defensible positions, not first-version launches",
        "Governance-first AI deployment — published frameworks, operational today",
    ]
    cw2 = Inches(5.9)
    for i, s in enumerate(shared):
        col = i % 2
        row = i // 2
        txt(slide, f"  ·  {s}", LM + col * (cw2 + Inches(0.4)), sy + row * Inches(0.4),
            cw2, Inches(0.36), size=10, color=C_TEXT)

    return slide


def s03_problem(prs):
    slide = blank_slide(prs)
    header(slide, "Corporate Travel Has a Missing Layer", 3)
    footer(slide, 3)
    ct = content_top()

    cw = Inches(5.6)
    cw2 = W - LM - RM - cw - Inches(0.4)
    divx = LM + cw + Inches(0.2)

    section_label(slide, "HOW DECISIONS ARE MADE TODAY", ct)
    steps = [
        ("01", "Policy set as a static document"),
        ("02", "Travellers book — often on consumer tools, outside policy"),
        ("03", "Finance reconciles spend after the fact"),
        ("04", "Carbon estimated annually from expense data, using averages"),
        ("05", "No real-time view of compliance or ESG alignment"),
    ]
    for i, (n, s) in enumerate(steps):
        y = ct + Inches(0.35) + i * Inches(0.72)
        rect(slide, LM, y + Inches(0.1), Inches(0.44), Inches(0.46), fill=C_NAVY)
        txt(slide, n, LM, y + Inches(0.1), Inches(0.44), Inches(0.48),
            size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(slide, s, LM + Inches(0.56), y + Inches(0.14), cw - Inches(0.6), Inches(0.42),
            size=11, color=C_TEXT)

    # Divider
    rect(slide, divx, ct, Inches(0.025), FOOTER_Y - ct - Inches(0.1), fill=C_BORDER)

    # Right col
    rx = divx + Inches(0.22)
    section_label(slide, "WHAT AI MAKES ACHIEVABLE NOW", ct)
    changes = [
        ("Policy compliance", "Assessed at booking time — not reconciled after the fact"),
        ("Carbon attribution", "Per-trip, route-level, using ICAO aircraft-type methodology"),
        ("Pricing context", "Historical and forward patterns surfaced at decision time"),
        ("Destination intelligence", "Risk, visa, and advisory data synthesised in real time"),
        ("Programme analytics", "Continuous visibility — not an annual report"),
    ]
    for i, (label, desc) in enumerate(changes):
        y = ct + Inches(0.35) + i * Inches(0.72)
        txt(slide, label, rx, y, CONTENT_W - cw - Inches(0.45), Inches(0.3),
            size=10.5, bold=True, color=C_NAVY)
        txt(slide, desc, rx, y + Inches(0.3), CONTENT_W - cw - Inches(0.45), Inches(0.38),
            size=10, color=C_TEXT)

    # Insight bar
    iy = FOOTER_Y - Inches(0.58)
    rect(slide, LM, iy, CONTENT_W, Inches(0.5), fill=C_LIGHT_BLUE_BG, line_rgb=C_BLUE, line_pt=0.5)
    txt(slide, "The opportunity is not to build a better booking tool. It is to build the intelligence infrastructure that makes every corporate travel decision better.",
        LM + Inches(0.15), iy + Inches(0.07), CONTENT_W - Inches(0.3), Inches(0.38),
        size=9.5, color=C_NAVY, italic=True)

    return slide


def s04_competitive(prs):
    slide = blank_slide(prs)
    header(slide, "Why Existing Platforms Fall Short", 4)
    footer(slide, 4)
    ct = content_top()

    txt(slide, "The gap is not from lack of trying. Navan, SAP Concur, TravelPerk, and Spotnana have all introduced intelligence features. The constraint is architectural, not commercial.",
        LM, ct, CONTENT_W, Inches(0.35), size=11, italic=True, color=C_MUTED)

    # Manual table (more control than add_table)
    ty = ct + Inches(0.45)
    col_ws = [Inches(2.1), Inches(3.1), Inches(6.6)]
    headers = ["Platform", "Design Priority", "Intelligence Limitation"]
    rows = [
        ["Navan", "Booking and expense consolidation",
         "Carbon tracking from spend categories. ICAO-level aircraft methodology not achievable from a booking-first data model."],
        ["SAP Concur", "ERP integration and expense management",
         "Sustainability module uses average industry factors. Policy compliance is rule-based matching, not AI reasoning."],
        ["TravelPerk", "SME booking experience",
         "Intelligence features are extensions to the booking workflow. Not designed for CSRD audit readiness."],
        ["Spotnana", "Modern booking infrastructure",
         "Platform layer. Intelligence is a third-party integration problem, not a core product capability."],
    ]

    row_h = Inches(0.42)
    header_h = Inches(0.42)
    all_rows = [headers] + rows

    for r, row in enumerate(all_rows):
        x = LM
        rh = header_h if r == 0 else row_h
        ry = ty + (header_h if r > 0 else 0) + (r - 1) * row_h if r > 0 else ty
        bg = C_NAVY if r == 0 else (C_LIGHT_BG if r % 2 == 0 else C_WHITE)
        for c, (cell, cw) in enumerate(zip(row, col_ws)):
            rect(slide, x, ry, cw, rh, fill=bg, line_rgb=C_BORDER, line_pt=0.4)
            tcol = C_WHITE if r == 0 else (C_NAVY if c == 0 else C_TEXT)
            txt(slide, cell, x + Inches(0.1), ry + Inches(0.07),
                cw - Inches(0.12), rh - Inches(0.1),
                size=9.5 if r == 0 else 10,
                bold=(r == 0 or c == 0), color=tcol, wrap=True)
            x += cw

    # Positioning box
    by = ty + header_h + len(rows) * row_h + Inches(0.18)
    rect(slide, LM, by, CONTENT_W, Inches(0.62), fill=C_NAVY)
    txt(slide, "DIS Travel is an intelligence platform that integrates with booking systems — not a booking platform with intelligence added. The architectural starting point determines what CSRD-compliant methodology is achievable.",
        LM + Inches(0.18), by + Inches(0.09), CONTENT_W - Inches(0.36), Inches(0.48),
        size=10.5, color=C_WHITE)

    txt(slide, "This architectural distinction is why the Skyscanner data partnership is structurally necessary, not interchangeable with other data sources.",
        LM, by + Inches(0.72), CONTENT_W, Inches(0.28),
        size=9.5, italic=True, color=C_MUTED)

    return slide


def s05_platform(prs):
    slide = blank_slide(prs)
    header(slide, "Introducing DIS Travel", 5)
    footer(slide, 5)
    ct = content_top()

    txt(slide, "An AI-Powered Travel Intelligence Platform for corporate buyers and travel-intensive organisations.",
        LM, ct, CONTENT_W, Inches(0.32), size=11, italic=True, color=C_MUTED)

    # Platform header bar
    ph = ct + Inches(0.42)
    rect(slide, LM, ph, CONTENT_W, Inches(0.5), fill=C_NAVY)
    txt(slide, "DIS TRAVEL  —  Travel Intelligence Platform",
        LM + Inches(0.2), ph + Inches(0.1), CONTENT_W - Inches(0.4), Inches(0.35),
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 4 capability cards
    caps = [
        ("POLICY\nINTELLIGENCE", "Compliance monitoring at booking time — not reconciled after"),
        ("DESTINATION\nINTELLIGENCE", "Risk, visa, health advisories — synthesised by AI at decision time"),
        ("CARBON\n& ESG", "Scope 3 per-trip attribution using ICAO / GHG Protocol methodology"),
        ("MARKET\nINTELLIGENCE", "Pricing trends and optimal booking window identification"),
    ]
    cw = Inches(2.93)
    cy = ph + Inches(0.5)
    ch = Inches(1.52)
    for i, (title, body) in enumerate(caps):
        x = LM + i * (cw + Inches(0.1))
        rect(slide, x, cy, cw, ch, fill=C_LIGHT_BG, line_rgb=C_BORDER, line_pt=0.75)
        txt(slide, title, x + Inches(0.12), cy + Inches(0.1), cw - Inches(0.24), Inches(0.6),
            size=9.5, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        txt(slide, body, x + Inches(0.12), cy + Inches(0.7), cw - Inches(0.24), Inches(0.72),
            size=9, color=C_TEXT, align=PP_ALIGN.CENTER)

    # AI layer
    ay = cy + ch + Inches(0.08)
    rect(slide, LM, ay, CONTENT_W, Inches(0.42), fill=C_LIGHT_BLUE_BG, line_rgb=C_BLUE, line_pt=0.75)
    txt(slide, "AI Intelligence & Analysis Layer",
        LM + Inches(0.15), ay + Inches(0.08), CONTENT_W - Inches(0.3), Inches(0.3),
        size=11, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # Data foundation
    dy = ay + Inches(0.42) + Inches(0.06)
    rect(slide, LM, dy, CONTENT_W, Inches(0.42), fill=RGBColor(240, 249, 255), line_rgb=RGBColor(147, 197, 253), line_pt=0.75)
    txt(slide, "Travel Content Data Foundation  ·  Skyscanner Partner Data  ·  ESG Data  ·  Other Sources",
        LM + Inches(0.15), dy + Inches(0.08), CONTENT_W - Inches(0.3), Inches(0.3),
        size=10, color=C_NAVY, align=PP_ALIGN.CENTER)

    # What it is not
    ny = dy + Inches(0.52)
    section_label(slide, "WHAT DIS TRAVEL IS NOT", ny)
    nots = [
        "Not a booking platform — does not compete with GDSs or OBTs",
        "Not a consumer product — built exclusively for enterprise buyers",
        "Not a TMC — does not replace the TMC relationship",
        "Not an expense tool — designed to integrate with Concur, SAP, etc.",
    ]
    cw2 = Inches(5.85)
    for i, n in enumerate(nots):
        col = i % 2
        row = i // 2
        txt(slide, f"  ✗  {n}", LM + col * (cw2 + Inches(0.42)), ny + Inches(0.32) + row * Inches(0.38),
            cw2, Inches(0.34), size=10, color=C_MUTED)

    return slide


def s06_differentiation(prs):
    slide = blank_slide(prs)
    header(slide, "Four Design Decisions That Define DIS Travel", 6)
    footer(slide, 6)
    ct = content_top()

    txt(slide, "These are architectural choices with specific commercial consequences — not positioning claims.",
        LM, ct, CONTENT_W, Inches(0.32), size=11, italic=True, color=C_MUTED)

    decisions = [
        ("01", "Intelligence Separable from Booking",
         "The intelligence layer operates independently of any specific booking system. DIS Travel integrates with an organisation's existing booking tool — it does not require replacing it. This removes the booking displacement objection from enterprise procurement entirely."),
        ("02", "Carbon as a First-Class Data Type",
         "Carbon attribution is calculated at trip-creation time, stored with full ICAO methodology provenance (inputs, methodology version, aircraft type, output), and queryable at any aggregation level. It is not a dashboard generated from expense exports — a distinction that CSRD audit readiness requires."),
        ("03", "GDPR as a Data Model Constraint",
         "Data minimisation and purpose limitation are enforced at the schema level — traveller PII is absent from the intelligence layer because the data model does not collect it there. This eliminates a class of compliance risk rather than managing it with policy documents."),
        ("04", "Governance-First AI Deployment",
         "AI components are not shipped until pre-deployment review is complete, risk class is assigned (Class 1–4 framework), and audit logging is confirmed. This adds deployment time. It is the correct tradeoff for a platform that enterprise procurement teams will scrutinise."),
    ]

    card_h = Inches(1.3)
    gap = Inches(0.08)
    for i, (num, title, body) in enumerate(decisions):
        y = ct + Inches(0.42) + i * (card_h + gap)
        rect(slide, LM, y, CONTENT_W, card_h, fill=C_LIGHT_BG, line_rgb=C_BORDER, line_pt=0.5)
        rect(slide, LM, y, Inches(0.06), card_h, fill=C_NAVY)
        rect(slide, LM + Inches(0.14), y + Inches(0.18), Inches(0.48), Inches(0.42), fill=C_NAVY)
        txt(slide, num, LM + Inches(0.14), y + Inches(0.18), Inches(0.48), Inches(0.44),
            size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(slide, title, LM + Inches(0.78), y + Inches(0.14), CONTENT_W - Inches(0.88), Inches(0.38),
            size=12, bold=True, color=C_NAVY)
        txt(slide, body, LM + Inches(0.78), y + Inches(0.52), CONTENT_W - Inches(0.88), Inches(0.68),
            size=10, color=C_TEXT)

    return slide


def s07_csrd(prs):
    slide = blank_slide(prs)
    header(slide, "CSRD Is Creating Urgent, Specific Demand", 7)
    footer(slide, 7)
    ct = content_top()

    cw = Inches(5.5)

    # Left column
    section_label(slide, "WHAT CSRD REQUIRES", ct)
    csrd = [
        ("In force from 2024", "Phased scope expansion through 2026"),
        ("Business travel = Scope 3 Category 6", "Mandatory measurement — estimation is insufficient"),
        ("Audit-defensible methodology required", "ICAO or equivalent standard expected by auditors"),
        ("Per-trip reporting capability", "Spend-category averages no longer adequate"),
        ("Tens of thousands of in-scope organisations", "Across the EU by full rollout"),
    ]
    for i, (bold, normal) in enumerate(csrd):
        y = ct + Inches(0.34) + i * Inches(0.74)
        rect(slide, LM, y + Inches(0.12), Inches(0.07), Inches(0.28), fill=C_BLUE)
        txt(slide, bold, LM + Inches(0.2), y + Inches(0.09), cw - Inches(0.24), Inches(0.3),
            size=11, bold=True, color=C_NAVY)
        txt(slide, normal, LM + Inches(0.2), y + Inches(0.38), cw - Inches(0.24), Inches(0.3),
            size=10, color=C_TEXT)

    # Divider
    dx = LM + cw + Inches(0.22)
    rect(slide, dx - Inches(0.22), ct, Inches(0.025), FOOTER_Y - ct - Inches(0.1), fill=C_BORDER)

    # Right column — comparison table
    rx = dx + Inches(0.0)
    rw = W - rx - RM - Inches(0.1)
    section_label(slide, "CURRENT STATE vs. DIS TRAVEL", ct)

    col_ws = [Inches(1.65), Inches(2.0), Inches(2.0)]
    headers = ["Dimension", "Today", "DIS Travel"]
    rows_data = [
        ["Timing", "Annual, retrospective", "Real-time, per booking"],
        ["Methodology", "Spend averages", "ICAO aircraft-type"],
        ["Granularity", "Category level", "Route + carrier level"],
        ["Audit readiness", "Estimated", "Methodology-logged"],
        ["Reporting", "Manual export", "Direct Scope 3 flow"],
    ]

    rh = Inches(0.44)
    all_rows = [headers] + rows_data
    for ri, row in enumerate(all_rows):
        x = rx
        y = ct + Inches(0.34) + ri * rh
        for ci, (cell, cw2) in enumerate(zip(row, col_ws)):
            bg = C_NAVY if ri == 0 else (C_LIGHT_BG if ri % 2 == 0 else C_WHITE)
            rect(slide, x, y, cw2, rh, fill=bg, line_rgb=C_BORDER, line_pt=0.4)
            col = C_WHITE if ri == 0 else (C_GREEN if ci == 2 and ri > 0 else (C_NAVY if ci == 0 else C_TEXT))
            txt(slide, cell, x + Inches(0.08), y + Inches(0.07), cw2 - Inches(0.1), rh - Inches(0.1),
                size=9.5, bold=(ri == 0 or ci == 0), color=col)
            x += cw2

    # Bottom note
    note_y = FOOTER_Y - Inches(0.48)
    txt(slide, "The accuracy of DIS Travel's carbon capability depends on route-level aircraft-type data — the specific data dependency that makes the Skyscanner partnership architecturally necessary.",
        LM, note_y, CONTENT_W, Inches(0.38), size=9.5, italic=True, color=C_MUTED)

    return slide


def s08_timing(prs):
    slide = blank_slide(prs)
    header(slide, "Three Forces Creating Demand Now", 8)
    footer(slide, 8)
    ct = content_top()

    txt(slide, "DIS Travel is positioned as the first platform to make CSRD-compliant travel intelligence architecturally achievable — a narrow, specific, and defensible position.",
        LM, ct, CONTENT_W, Inches(0.35), size=11, italic=True, color=C_MUTED)

    forces = [
        ("01", "CSRD Compliance\nIs Active and Urgent",
         "The EU Corporate Sustainability Reporting Directive is in force. Companies in scope are actively building Scope 3 measurement infrastructure. Business travel is a primary category. Compliance deadlines create procurement budget and urgency that did not exist two years ago.",
         C_NAVY),
        ("02", "AI Infrastructure Has\nReached Enterprise Viability",
         "The large language models, data pipelines, and inference infrastructure that DIS Travel requires are now commercially available at the cost, latency, and reliability thresholds that enterprise software demands. Building this was technically possible but economically marginal three years ago.",
         C_TEAL),
        ("03", "Corporate Travel Buyers Are\nRedefining Requirements",
         "Corporate travel budgets have returned to pre-2020 levels — but the conversations around them have changed. CFOs are asking about ESG alignment. Procurement teams are asking about audit readiness. Sustainability teams are asking about Scope 3 methodology. The tool set has not kept pace.",
         C_PURPLE),
    ]

    cw = Inches(3.88)
    gap = Inches(0.18)

    for i, (num, title, body, col) in enumerate(forces):
        x = LM + i * (cw + gap)
        num_h = Inches(0.52)
        title_h = Inches(0.88)
        body_h = Inches(3.2)

        rect(slide, x, ct + Inches(0.44), cw, num_h, fill=col)
        txt(slide, num, x + Inches(0.15), ct + Inches(0.5), cw - Inches(0.3), Inches(0.42),
            size=22, bold=True, color=C_WHITE)

        rect(slide, x, ct + Inches(0.44) + num_h, cw, title_h, fill=RGBColor(246, 248, 250), line_rgb=C_BORDER, line_pt=0.5)
        txt(slide, title, x + Inches(0.15), ct + Inches(0.6) + num_h, cw - Inches(0.3), title_h - Inches(0.2),
            size=12, bold=True, color=col)

        rect(slide, x, ct + Inches(0.44) + num_h + title_h, cw, body_h, fill=C_WHITE, line_rgb=C_BORDER, line_pt=0.5)
        txt(slide, body, x + Inches(0.15), ct + Inches(0.62) + num_h + title_h, cw - Inches(0.3), body_h - Inches(0.24),
            size=10.5, color=C_TEXT)

    return slide


def s09_skyscanner(prs):
    slide = blank_slide(prs)
    header(slide, "Why Skyscanner Is Architecturally Necessary", 9)
    footer(slide, 9)
    ct = content_top()

    txt(slide, "We have evaluated the data landscape carefully. Skyscanner's value to DIS Travel is specific, not generic.",
        LM, ct, CONTENT_W, Inches(0.32), size=11, italic=True, color=C_MUTED)

    # Table
    col_ws = [Inches(3.8), Inches(7.95)]
    headers = ["What DIS Travel Requires", "Why Skyscanner — Specifically"]
    rows_data = [
        ["Route-level aircraft metadata",
         "Required for ICAO/GHG Protocol carbon attribution. Available at depth through Skyscanner's partner programme in a way consumer-facing sources do not provide."],
        ["European route coverage",
         "Skyscanner's depth in European routes aligns with DIS Travel's launch market. Accuracy across these routes is critical for initial enterprise customers."],
        ["Structured B2B partner infrastructure",
         "Skyscanner has built a partner programme with the data contracts, SLAs, and API design that enterprise integrations require. Reduces integration risk for both parties."],
        ["Real-time pricing signals",
         "Market intelligence requires live data. Skyscanner's real-time pricing infrastructure is the signal layer that makes booking-window recommendations actionable."],
    ]

    rh = Inches(0.68)
    all_rows = [headers] + rows_data
    ty = ct + Inches(0.42)
    for ri, row in enumerate(all_rows):
        x = LM
        y = ty + ri * (Inches(0.44) if ri == 0 else rh) if ri == 0 else ty + Inches(0.44) + (ri - 1) * rh
        rh2 = Inches(0.44) if ri == 0 else rh
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            bg = C_NAVY if ri == 0 else (C_LIGHT_BG if ri % 2 == 1 else C_WHITE)
            rect(slide, x, y, cw, rh2, fill=bg, line_rgb=C_BORDER, line_pt=0.4)
            col = C_WHITE if ri == 0 else (C_NAVY if ci == 0 else C_TEXT)
            txt(slide, cell, x + Inches(0.12), y + Inches(0.06), cw - Inches(0.16), rh2 - Inches(0.08),
                size=9.5 if ri == 0 else 10, bold=(ri == 0 or ci == 0), color=col, wrap=True)
            x += cw

    # What DIS Travel brings
    by = ty + Inches(0.44) + len(rows_data) * rh + Inches(0.14)
    section_label(slide, "WHAT DIS TRAVEL BRINGS TO SKYSCANNER", by)
    brings = [
        "Access to corporate buyers — a B2B market Skyscanner's consumer product does not directly serve",
        "CSRD-driven demand: enterprise buyers with compliance obligations and dedicated budget for intelligence tools",
        "Long-term architectural investment — Skyscanner as the foundation data layer for a decade-horizon platform",
        "Distribution: when DIS Travel identifies the optimal option, the booking connects to Skyscanner's supplier network",
    ]
    bw = Inches(5.85)
    for i, b in enumerate(brings):
        col = i % 2
        row = i // 2
        txt(slide, f"  ·  {b}", LM + col * (bw + Inches(0.42)), by + Inches(0.3) + row * Inches(0.4),
            bw, Inches(0.36), size=10, color=C_TEXT)

    # Core positioning
    py = FOOTER_Y - Inches(0.52)
    rect(slide, LM, py, CONTENT_W, Inches(0.44), fill=C_LIGHT_BLUE_BG, line_rgb=C_BLUE, line_pt=0.5)
    txt(slide, "Skyscanner serves travellers.  DIS Travel serves the organisations that employ them.  The relationship is complementary, not competitive.",
        LM + Inches(0.15), py + Inches(0.08), CONTENT_W - Inches(0.3), Inches(0.32),
        size=10.5, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    return slide


def s10_partnership(prs):
    slide = blank_slide(prs)
    header(slide, "What We Are Proposing to Build Together", 10)
    footer(slide, 10)
    ct = content_top()

    # 3 component cards
    cw = Inches(3.88)
    ch = Inches(2.58)
    gap = Inches(0.18)
    cards = [
        ("01  DATA FOUNDATION", C_NAVY,
         "Structured access to Skyscanner's travel content as the foundation layer for DIS Travel's intelligence:\n\n·  Route and schedule data\n·  Real-time pricing and availability\n·  Aircraft type per route (ICAO carbon)\n·  Historical fare patterns"),
        ("02  DISTRIBUTION ALIGNMENT", C_TEAL,
         "When DIS Travel intelligence drives a booking decision, the booking connects to Skyscanner's supplier network.\n\nCorporate bookings via intelligence tools represent high-intent, context-rich transactions — distinct from consumer metasearch traffic."),
        ("03  CARBON DATA COLLABORATION", C_PURPLE,
         "Co-development of route-level carbon intelligence anchored in Skyscanner's route and aircraft data, validated against ICAO/GHG Protocol methodology.\n\nNeither organisation can build this at the same accuracy independently."),
    ]
    for i, (title, col, body) in enumerate(cards):
        x = LM + i * (cw + gap)
        rect(slide, x, ct + Inches(0.06), cw, Inches(0.46), fill=col)
        txt(slide, title, x + Inches(0.12), ct + Inches(0.12), cw - Inches(0.24), Inches(0.36),
            size=10, bold=True, color=C_WHITE)
        rect(slide, x, ct + Inches(0.52), cw, ch - Inches(0.46), fill=C_LIGHT_BG, line_rgb=C_BORDER, line_pt=0.5)
        txt(slide, body, x + Inches(0.14), ct + Inches(0.64), cw - Inches(0.28), ch - Inches(0.7),
            size=9.5, color=C_TEXT)

    # Timeline
    ty = ct + ch + Inches(0.28)
    section_label(slide, "INDICATIVE PARTNERSHIP TIMELINE", ty)
    ty += Inches(0.3)

    tl_data = [
        ["Phase", "Period", "Milestone"],
        ["Exploratory", "Q3–Q4 2026", "Leadership and technical conversations; assess strategic fit"],
        ["Partnership Design", "Q4 2026 – Q1 2027", "Commercial structure, data agreement, technical specification"],
        ["Integration Build", "Q1–Q2 2027", "Data connection, integration testing, quality assurance"],
        ["Enterprise Pilot", "Q2–Q3 2027", "First enterprise pilots with DIS Travel + Skyscanner content"],
        ["Scale", "2027–2028", "Expand enterprise accounts; deepen integration"],
    ]
    col_ws = [Inches(2.4), Inches(2.5), Inches(7.0)]
    rh = Inches(0.43)
    for ri, row in enumerate(tl_data):
        x = LM
        y = ty + ri * rh
        for ci, (cell, cw2) in enumerate(zip(row, col_ws)):
            bg = C_NAVY if ri == 0 else (C_LIGHT_BG if ri % 2 == 0 else C_WHITE)
            rect(slide, x, y, cw2, rh, fill=bg, line_rgb=C_BORDER, line_pt=0.4)
            col = C_WHITE if ri == 0 else (C_NAVY if ci == 0 else C_TEXT)
            txt(slide, cell, x + Inches(0.1), y + Inches(0.07), cw2 - Inches(0.12), rh - Inches(0.1),
                size=9.5, bold=(ri == 0 or ci == 0), color=col)
            x += cw2

    txt(slide, "Timeline is indicative — subject to partnership discussion. We are not proposing Skyscanner commit to anything before the strategic fit is established.",
        LM, FOOTER_Y - Inches(0.35), CONTENT_W, Inches(0.28),
        size=8.5, italic=True, color=C_MUTED)

    return slide


def s11_why_dis(prs):
    slide = blank_slide(prs)
    header(slide, "Why DIS Group as a Long-Term Partner", 11)
    footer(slide, 11)
    ct = content_top()

    txt(slide, "The question enterprise partners ask is not whether the product idea is interesting. It is whether the team behind it can execute.",
        LM, ct, CONTENT_W, Inches(0.35), size=11, italic=True, color=C_MUTED)

    cw = Inches(3.88)
    ch = Inches(5.1)
    gap = Inches(0.18)
    cols = [
        ("OPERATING FOUNDATION", C_NAVY,
         "DIS Group is not a startup built around a founding pitch.\n\nWe operate three businesses today — in recycling operations, energy trading intelligence, and ESG data infrastructure.\n\nThe governance frameworks and technology standards governing DIS Travel are adapted from frameworks operational in those businesses.\n\nDIS Travel is the application of a proven development discipline to a new market — not a first attempt at building something complex."),
        ("GOVERNANCE YOU CAN AUDIT", C_TEAL,
         "Our technology governance is operational today:\n\n·  AI governance policy — four-class risk classification, pre-deployment review, incident response\n\n·  Architecture decisions via ADRs — documented rationale for every major choice\n\n·  GDPR compliance — EU-incorporated, EU data residency capability, enterprise DPA-ready\n\n·  Published documentation — available to share, not just to describe"),
        ("TRANSPARENCY AS PRACTICE", C_PURPLE,
         "Every document in this package is honest about where we are.\n\nOur development status table shows clearly what is complete, in progress, and designed-but-not-built.\n\nWe do not claim customers we do not have or revenues we have not earned.\n\nThe most enduring technology partnerships are built on accurate mutual understanding — and that is the basis on which we are approaching this conversation."),
    ]
    for i, (title, col, body) in enumerate(cols):
        x = LM + i * (cw + gap)
        rect(slide, x, ct + Inches(0.44), cw, Inches(0.44), fill=col)
        txt(slide, title, x + Inches(0.14), ct + Inches(0.5), cw - Inches(0.28), Inches(0.34),
            size=10, bold=True, color=C_WHITE)
        rect(slide, x, ct + Inches(0.88), cw, ch - Inches(0.44), fill=C_LIGHT_BG, line_rgb=C_BORDER, line_pt=0.5)
        txt(slide, body, x + Inches(0.16), ct + Inches(1.02), cw - Inches(0.32), ch - Inches(0.7),
            size=9.5, color=C_TEXT)

    return slide


def s12_next_steps(prs):
    slide = blank_slide(prs)
    header(slide, "A Structured Path to Partnership", 12)
    footer(slide, 12)
    ct = content_top()

    txt(slide, "What we are asking for is a conversation. Not a pilot. Not a data access agreement. Not a commercial commitment.",
        LM, ct, CONTENT_W, Inches(0.35), size=12, bold=True, color=C_NAVY)

    section_label(slide, "WHAT SKYSCANNER COMMITS AT EACH STAGE", ct + Inches(0.44))
    risk_data = [
        ["Stage", "Skyscanner Commits", "Skyscanner Receives"],
        ["Initial conversation\n(45 minutes)", "One meeting\nNo commercial or data commitment", "A structured assessment of whether strategic fit warrants further investment"],
        ["Partnership design\n(60–90 days)", "Dedicated time from partnerships\nand technical teams", "Fully specified commercial proposal, technical integration design, and data agreement draft"],
        ["Integration build", "Defined data access under\nsigned agreement", "Corporate buyer channel and CSRD-driven distribution Skyscanner's consumer product does not currently reach"],
    ]
    col_ws = [Inches(2.5), Inches(3.7), Inches(5.6)]
    rh = Inches(0.72)
    ty = ct + Inches(0.76)
    for ri, row in enumerate(risk_data):
        x = LM
        rh2 = Inches(0.44) if ri == 0 else rh
        y = ty if ri == 0 else ty + Inches(0.44) + (ri - 1) * rh
        for ci, (cell, cw) in enumerate(zip(row, col_ws)):
            bg = C_NAVY if ri == 0 else (C_LIGHT_BG if ri % 2 == 1 else C_WHITE)
            rect(slide, x, y, cw, rh2, fill=bg, line_rgb=C_BORDER, line_pt=0.4)
            col = C_WHITE if ri == 0 else (C_NAVY if ci == 0 else C_TEXT)
            txt(slide, cell, x + Inches(0.12), y + Inches(0.06), cw - Inches(0.16), rh2 - Inches(0.08),
                size=10, bold=(ri == 0 or ci == 0), color=col, wrap=True)
            x += cw

    # What we bring
    by = ty + Inches(0.44) + len(risk_data[1:]) * rh + Inches(0.18)
    section_label(slide, "WHAT WE BRING TO THE INITIAL CONVERSATION", by)
    bring_items = [
        ("Technical depth", "Architecture walkthrough and specific data dependencies"),
        ("Commercial clarity", "Value exchange articulated precisely; open to Skyscanner's preferred commercial structure"),
        ("Honesty about timelines", "We will not overstate build progress or compress realistic partnership timelines"),
    ]
    bw = (CONTENT_W - Inches(0.36)) / 3
    for i, (title, body) in enumerate(bring_items):
        x = LM + i * (bw + Inches(0.18))
        by2 = by + Inches(0.3)
        rect(slide, x, by2, bw, Inches(0.38), fill=C_NAVY)
        txt(slide, title, x + Inches(0.12), by2 + Inches(0.06), bw - Inches(0.2), Inches(0.28),
            size=10, bold=True, color=C_WHITE)
        rect(slide, x, by2 + Inches(0.38), bw, Inches(0.78), fill=C_LIGHT_BG, line_rgb=C_BORDER, line_pt=0.5)
        txt(slide, body, x + Inches(0.12), by2 + Inches(0.48), bw - Inches(0.2), Inches(0.6),
            size=9.5, color=C_TEXT)

    # CTA
    cta_y = FOOTER_Y - Inches(0.58)
    rect(slide, LM, cta_y, CONTENT_W, Inches(0.5), fill=C_NAVY)
    txt(slide, "Proposed next step: a 45-minute call at a time that suits your team.   Contact: [ceo@disgroup.ie]",
        LM + Inches(0.2), cta_y + Inches(0.1), CONTENT_W - Inches(0.4), Inches(0.34),
        size=11.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    return slide


# ─── MAIN ───────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    out = "/Users/user1/Desktop/DIS-Foundation/partners/skyscanner/final/05_DIS_Travel_Presentation.pptx"
    os.makedirs(os.path.dirname(out), exist_ok=True)

    prs = new_prs()
    s01_title(prs)
    s02_about(prs)
    s03_problem(prs)
    s04_competitive(prs)
    s05_platform(prs)
    s06_differentiation(prs)
    s07_csrd(prs)
    s08_timing(prs)
    s09_skyscanner(prs)
    s10_partnership(prs)
    s11_why_dis(prs)
    s12_next_steps(prs)

    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")
